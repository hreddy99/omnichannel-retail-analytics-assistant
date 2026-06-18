"""
Builds the capstone presentation deck (16:9) for the Omnichannel Retail
Analytics Assistant. Run: python presentation/build_deck.py
Produces presentation/Omnichannel_Retail_Analytics_Assistant.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- palette (Carnegie Mellon red + app accents) ----
CMU_RED = RGBColor(0xC4, 0x12, 0x30)
SLATE = RGBColor(0x1E, 0x29, 0x3B)
GREY = RGBColor(0x47, 0x55, 0x69)
LIGHT = RGBColor(0xF1, 0xF5, 0xF9)
BLUE = RGBColor(0x25, 0x63, 0xEB)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Segoe UI"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def _box(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    return tb


def _fill(slide, l, t, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def _para(tf, text, size, color, bold=False, bullet=False, level=0, space_after=8, italic=False):
    p = tf.paragraphs[0] if (tf.paragraphs and tf.paragraphs[0].text == "" and not tf._txBody.findall(
        '{http://schemas.openxmlformats.org/drawingml/2006/main}p')[0].findall(
        '{http://schemas.openxmlformats.org/drawingml/2006/main}r')) else tf.add_paragraph()
    p.text = text
    p.level = level
    p.space_after = Pt(space_after)
    for r in p.runs:
        r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = color; r.font.name = FONT
    return p


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def content_slide(number, kicker, title, bullets, note, accent=CMU_RED):
    """bullets: list of (text, level, bold) ; level 0 main, 1 sub."""
    s = prs.slides.add_slide(BLANK)
    _fill(s, 0, 0, SW, Inches(0.18), accent)            # top rule
    # kicker + title
    k = _box(s, Inches(0.6), Inches(0.45), Inches(12.1), Inches(0.4))
    _para(k.text_frame, kicker.upper(), 13, accent, bold=True, space_after=2)
    tt = _box(s, Inches(0.6), Inches(0.85), Inches(12.1), Inches(1.0))
    _para(tt.text_frame, title, 30, SLATE, bold=True, space_after=0)
    # body
    body = _box(s, Inches(0.7), Inches(1.95), Inches(11.9), Inches(5.0))
    tf = body.text_frame
    first = True
    for item in bullets:
        text, level, bold = item
        color = SLATE if level == 0 else GREY
        size = 19 if level == 0 else 16
        p = _para(tf, ("•  " if level == 0 else "–  ") + text, size, color,
                  bold=bold, level=level, space_after=9 if level == 0 else 5)
        if first:
            first = False
    # footer
    f = _box(s, Inches(0.6), Inches(7.0), Inches(12.1), Inches(0.4))
    fp = _para(f.text_frame, f"Omnichannel Retail Analytics Assistant  ·  CMU Capstone", 10, GREY)
    pg = _box(s, Inches(12.4), Inches(7.0), Inches(0.6), Inches(0.4))
    _para(pg.text_frame, str(number), 11, GREY, bold=True)
    notes(s, note)
    return s


# ============================================================= SLIDE 1 — Title
s = prs.slides.add_slide(BLANK)
_fill(s, 0, 0, SW, SH, SLATE)
_fill(s, 0, Inches(2.55), SW, Inches(0.10), CMU_RED)
t = _box(s, Inches(0.9), Inches(1.4), Inches(11.5), Inches(1.2))
_para(t.text_frame, "Omnichannel Retail Analytics Assistant", 44, WHITE, bold=True)
sub = _box(s, Inches(0.9), Inches(2.75), Inches(11.5), Inches(1.2))
_para(sub.text_frame, "A governed, multi-agent AI analyst that explains business performance —",
      22, LIGHT, space_after=2)
_para(sub.text_frame, "grounded in certified metrics, read-only, and free to run locally.",
      22, LIGHT)
meta = _box(s, Inches(0.9), Inches(5.6), Inches(11.5), Inches(1.4))
_para(meta.text_frame, "Final Capstone Presentation  ·  Carnegie Mellon University", 16, WHITE, bold=True, space_after=4)
_para(meta.text_frame, "Presented by Sreddy  ·  June 2026", 14, LIGHT)
notes(s, "TIME 0:00–0:45. Hook: 'When digital conversion drops, a retail leader asks one "
        "question — why? Answering it today means pinging five teams and waiting two days. "
        "I built an AI analyst that answers it in seconds — but, crucially, only from "
        "governed, certified data, and it never guesses.' State your name, that this is a "
        "working prototype you'll demo live. Set expectations: 20 minutes, ~7 of it live demo.")

# ============================================================= SLIDE 2 — Problem
content_slide(
    2, "The problem", "Cross-functional 'why' questions are slow, and BI tools guess",
    [("A KPI moves — conversion, revenue, returns — and the answer spans many teams at once: "
      "marketing, merchandising, fulfillment, digital, service, finance, vendors.", 0, False),
     ("Today: analysts hand-pull data from each silo; it takes days and answers conflict.", 0, False),
     ("Generic 'chat-with-your-data' tools are fast but ungoverned — they hallucinate metric "
      "definitions, query whatever they want, and assert causation from correlation.", 0, False),
     ("In retail, a confidently wrong answer routed to the wrong team is worse than a slow one.", 0, True)],
    "TIME 0:45–2:30. Make it concrete and personal to retail. The pain is twofold: (1) SPEED — "
    "cross-domain questions are slow because the data is siloed; (2) TRUST — the fast AI tools on "
    "the market are ungoverned: they invent definitions, run arbitrary SQL, and claim causation. "
    "Land the last line hard: in an enterprise, an authoritative wrong answer is a liability. "
    "That tension — fast AND trustworthy — is the whole project.")

# ============================================================= SLIDE 3 — Users
content_slide(
    3, "Who it's for", "Built around real retail roles and the questions they ask",
    [("Leadership / Analytics — 'Why did digital conversion drop yesterday vs the prior week?'", 0, False),
     ("Marketing — 'Did paid-social traffic convert below baseline?'", 0, False),
     ("Merchandising — 'Which categories are stocking out? Which vendor do we alert?'", 0, False),
     ("Fulfillment — 'Where are delivery delays worst, and is it driving cancellations?'", 0, False),
     ("Customer Service & Finance — 'Did contacts spike? Why doesn't net revenue match gross?'", 0, False),
     ("One assistant, one governed source of truth — every team gets a consistent answer.", 0, True)],
    "TIME 2:30–3:45. These are the personas in my catalog — each maps to a specialized analyst "
    "agent and an owner. The key design idea: every role asks its own question in plain English, "
    "but they all resolve against the SAME certified metric definitions. No more 'my conversion "
    "number doesn't match yours.' That consistency is a governance feature, not an accident.")

# ============================================================= SLIDE 4 — What I built
content_slide(
    4, "The solution", "A governed multi-agent analyst — answer + evidence + owner action",
    [("Ask in natural language → the system retrieves CERTIFIED definitions, runs read-only "
      "queries, reasons across domains, and returns a grounded answer.", 0, False),
     ("Every answer ships with: the evidence (the actual numbers + chart), the certified "
      "definition used, a confidence label, caveats, and an owner-routed next action.", 0, False),
     ("Three answer modes:", 0, True),
     ("Flagship investigation — 'why did conversion drop?' fans out a 7-analyst team.", 1, False),
     ("Direct analytics & themed reviews — 15 + 12 governed questions across all domains.", 1, False),
     ("Executive briefing — ranks the biggest cross-functional issues and routes each to an owner.", 1, False),
     ("100% free & local: synthetic data, local models, no paid APIs, nothing leaves the laptop.", 0, True)],
    "TIME 3:45–5:30. This is the 'what' before the 'how.' Emphasize the answer is never just a "
    "number — it's answer + evidence + definition + confidence + caveat + owner action. That "
    "package is what makes it enterprise-grade. Mention the three modes briefly (you'll demo two). "
    "Close on feasibility: the plan claimed a free, local, read-only MVP runnable on a personal "
    "PC — this prototype proves it end to end. That matters for a capstone: it actually runs.")

# ============================================================= SLIDE 5 — Architecture
content_slide(
    5, "Architecture", "Governed knowledge + a multi-agent reasoning loop on a modern data stack",
    [("Data platform (medallion): synthetic Bronze → Silver fact/dim tables → Gold certified "
      "metrics, queried read-only in DuckDB (18 governed tables).", 0, False),
     ("Governed knowledge layer — the source of truth:", 0, True),
     ("YAML catalog: metrics, tables, drivers, guardrails, owners (versioned + hashed).", 1, False),
     ("ChromaDB vector index for retrieval; NetworkX graph for metric→driver→table→owner.", 1, False),
     ("Reasoning layer: a LangGraph workflow orchestrates a team of specialized agents with a "
      "Critic, beam search, and a human-review gate.", 0, False),
     ("Presentation: a Streamlit app with four trust levels (answer → evidence → trust → audit).", 0, False),
     ("Optional local LLM (Ollama) only DRAFTS prose — it never invents numbers or definitions.", 0, True)],
    "TIME 5:30–7:30. Walk top-to-bottom. Stress the architectural principle: the LLM is "
    "deliberately on a short leash — it only phrases the final sentence; all facts come from "
    "governed SQL over certified metrics. If Ollama isn't running, a deterministic fallback "
    "writes the prose, so the demo never depends on a model download. The medallion framing shows "
    "I aligned to a real enterprise data platform, not a toy script. Point at the YAML catalog as "
    "'source of truth' — that one decision drives the whole trust story.")

# ============================================================= SLIDE 6 — Governance / trust
content_slide(
    6, "Trust by construction", "Governance is enforced in code, not promised in a prompt",
    [("Read-only guardrail: a SQL validator allows only SELECT/WITH over APPROVED tables — "
      "any write or unapproved table is blocked before execution.", 0, False),
     ("Certified definitions only: metrics resolve from the versioned YAML catalog; retrieved "
      "context is validated against it before use.", 0, False),
     ("Causality is LABELED, never asserted — 'likely driver' / 'possible contributor', not 'caused'.", 0, True),
     ("Human-in-the-loop: every recommendation routes to an owner for review; no system is ever "
      "modified automatically.", 0, False),
     ("Full observability: catalog version + content hash, freshness gate, and a step-by-step "
      "audit log behind every answer.", 0, False)],
    "TIME 7:30–9:00. This is the heart of a CMU-grade project — responsible AI made concrete. "
    "Hit each as an ENFORCED control: the guardrail isn't advice to the model, it's a validator "
    "that refuses the query. Causality labeling is the intellectually honest move — with "
    "observational data you cannot prove cause, so the system is explicit about it. Human-in-the-"
    "loop + read-only means it's safe to deploy in a governed enterprise. Tell them you'll SHOW "
    "the guardrail refuse a write during the demo.")

# ============================================================= SLIDE 7 — Multi-agent loop
content_slide(
    7, "Multi-agent reasoning", "Why a team of agents — and when NOT to use one",
    [("A conversion drop can have several simultaneous causes across domains — so the Orchestrator "
      "dispatches 7 specialized analysts IN PARALLEL (one read-only query each).", 0, False),
     ("Pipeline: classify → retrieve → validate → graph-relate → baseline → ToT gate → dispatch "
      "team → Critic scores each (0–14 rubric) → beam search → evidence gate → human review → synthesize.", 0, False),
     ("The Critic prunes ungoverned hypotheses; beam search keeps the strongest drivers; "
      "corroborating signals (service/finance/vendor) are reported separately, not as causes.", 0, False),
     ("Deliberate, not multi-agent for its own sake:", 0, True),
     ("Used only when the question is cross-domain; a narrow question uses one analyst.", 1, False),
     ("A failing/slow agent is isolated and excluded — it degrades gracefully, never sinks the team.", 1, False)],
    "TIME 9:00–11:00. Justify the multi-agent design as an engineering decision with trade-offs — "
    "examiners reward this. WHY multi-agent: the problem is genuinely parallel and specialized. "
    "WHEN NOT: narrow questions take a single analyst; the ToT gate decides. Mention the 0–14 "
    "Critic rubric (metric validated, approved graph path, SQL safety, evidence strength, "
    "freshness, business relevance, caveats) — that's how a hypothesis earns 'likely driver.' "
    "Name the trade-offs you accepted (coordination overhead, complexity, non-determinism) and "
    "how you bounded each (thread pool + timeout, one shared contract, deterministic re-sort).")

# ============================================================= SLIDE 8 — DEMO
s = prs.slides.add_slide(BLANK)
_fill(s, 0, 0, SW, SH, SLATE)
_fill(s, 0, Inches(2.7), SW, Inches(0.10), CMU_RED)
k = _box(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(0.5))
_para(k.text_frame, "LIVE DEMO", 18, CMU_RED, bold=True)
t = _box(s, Inches(0.9), Inches(2.9), Inches(11.5), Inches(1.0))
_para(t.text_frame, "Seeing it run on a seeded 'bad day'", 36, WHITE, bold=True)
b = _box(s, Inches(0.9), Inches(4.1), Inches(11.5), Inches(2.3))
tf = b.text_frame
_para(tf, "1.  Flagship: 'Why did digital conversion drop yesterday?' — watch the team dispatch, "
          "the live trace, drivers + evidence chart, and the owner-routed actions.", 17, LIGHT, space_after=10)
_para(tf, "2.  Guardrail: ask it to 'update the paid-social budget' — it refuses (read-only).", 17, LIGHT, space_after=10)
_para(tf, "3.  Executive briefing: the multi-agent path ranks the biggest cross-functional issues.", 17, LIGHT)
notes(s, "TIME 11:00–18:00 (your 7-minute centerpiece — see DEMO_TALKING_POINTS.md for the "
        "click-by-click script). Keep the app already running before you start (python tools.py "
        "run). Narrate intent, not mechanics: say what a user wants, then show the system earn the "
        "answer. Order: (1) run the conversion investigation; while the live trace streams, narrate "
        "the governed steps; land on the Business answer, then the Evidence chart, then the "
        "owner actions + human-review banner. (2) Show the read-only guardrail refusing a write — "
        "this always gets a reaction. (3) Run one Executive briefing to show the cross-functional "
        "ranking. If anything is slow, talk over it — you have the trace to point at. Have a backup "
        "screenshot/GIF in case the laptop misbehaves on the recording.")

# ============================================================= SLIDE 9 — Evaluation
content_slide(
    9, "Does it work?", "Evaluation: reproducible, validated, and exception-free",
    [("Seeded 'bad day': a fixed-seed synthetic world with a coherent, known story — conversion "
      "down ~23% driven by paid-social shift, stockouts, regional delays, a funnel dip.", 0, False),
     ("Automated validation: 15 data/guardrail/routing checks all pass — the seeded signals are "
      "present, writes are blocked, and every question routes to the right path.", 0, False),
     ("Coverage: 1 flagship investigation + 3 executive briefings + 15 analytics + 12 themed "
      "reviews — all run with real numbers and a matching chart; UI renders with 0 exceptions.", 0, False),
     ("Reproducible by design: same seed → same evidence, so results are inspectable and gradeable.", 0, True)],
    "TIME 18:00–19:00. Evaluation is where capstones win or lose. Two layers: (1) the data has a "
    "KNOWN ground truth because I seeded it, so I can assert the system finds the right drivers; "
    "(2) automated checks gate the whole thing — data validity, guardrail enforcement, and "
    "question-routing accuracy. Everything is reproducible from a fixed seed, which is exactly "
    "what makes an AI system auditable. Mention 0 render exceptions across all pages as a quality bar.")

# ============================================================= SLIDE 10 — Trade-offs / decisions
content_slide(
    10, "Engineering judgment", "Key decisions & honest trade-offs",
    [("Free/local stack (DuckDB, ChromaDB, optional Ollama) — zero cost and private, at the price "
      "of synthetic data and smaller models. Right call for a governed prototype.", 0, False),
     ("LLM drafts prose only — gives natural language without risking hallucinated facts.", 0, False),
     ("Synthetic data — lets me seed a known story and avoid PII, but magnitudes are illustrative.", 0, False),
     ("Bounded reasoning — query budget, beam width, and depth cap cost and keep latency predictable.", 0, False),
     ("Limitations I'm explicit about: observational data (no proven causation), single-day "
      "scenario, and a catalog scoped to the demo domains.", 0, True)],
    "TIME 19:00–19:40. Show you made choices like an engineer and you know the limits. Frame each "
    "trade-off as deliberate with a clear rationale, then own the limitations without prompting — "
    "examiners trust a presenter who names their own weaknesses. The honesty here pairs with the "
    "causality-labeling slide: this is a system that knows what it doesn't know.")

# ============================================================= SLIDE 11 — Impact / future / learnings
content_slide(
    11, "Impact & what's next", "From a 2-day fire drill to a governed answer in seconds",
    [("Impact: collapses a cross-team investigation into one grounded, owner-routed answer — "
      "fast AND trustworthy, the tension we started with.", 0, False),
     ("What I learned: governance is an architecture problem, not a prompt; the hard part of "
      "agentic AI is restraint, evaluation, and observability — not generating text.", 0, False),
     ("Next steps:", 0, True),
     ("Connect to real warehouse tables behind the same YAML governance layer.", 1, False),
     ("Add multi-day trends and statistical attribution to strengthen causal claims.", 1, False),
     ("Expand the analyst roster and let owners action recommendations through approval workflows.", 1, False)],
    "TIME 19:40–20:00. Bring it home to the opening promise: fast AND trustworthy. State the one "
    "big lesson crisply — 'the hard part of agentic AI isn't generating answers, it's governing, "
    "evaluating, and observing them.' Give a credible roadmap that shows the architecture extends "
    "to production (swap synthetic for real tables behind the same catalog). Then thank them and "
    "invite questions.")

# ============================================================= SLIDE 12 — Thank you
s = prs.slides.add_slide(BLANK)
_fill(s, 0, 0, SW, SH, SLATE)
_fill(s, 0, Inches(3.4), SW, Inches(0.10), CMU_RED)
t = _box(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(1.0))
_para(t.text_frame, "Thank you", 48, WHITE, bold=True)
sub = _box(s, Inches(0.9), Inches(3.7), Inches(11.5), Inches(1.2))
_para(sub.text_frame, "Questions?", 24, LIGHT, space_after=6)
_para(sub.text_frame, "Governed, multi-agent, read-only, free to run locally — and it actually runs.",
      16, LIGHT, italic=True)
notes(s, "Q&A. Likely questions + crisp answers: (1) 'How do you know the drivers are real?' — "
        "seeded ground truth + the Critic's evidence rubric; I label likelihood, not proof. "
        "(2) 'Why not just GPT-4 over the warehouse?' — governance: certified definitions, "
        "read-only validator, audit trail; an ungoverned LLM can't give consistent, safe answers. "
        "(3) 'Is it production-ready?' — the governance architecture is; swap synthetic data for "
        "real tables behind the same catalog. (4) 'Why multi-agent?' — parallel, specialized, "
        "cross-domain; single-domain questions use one analyst. Keep answers to 20–30 seconds.")

import os
out = os.path.join(os.path.dirname(__file__), "Omnichannel_Retail_Analytics_Assistant.pptx")
prs.save(out)
print("Saved:", out, "| slides:", len(prs.slides._sldIdLst))
