"""
8-minute capstone deck — designed, image-composed slides (verifiable) with a drawn
architecture diagram, assembled into a PPTX with speaker notes.
Flow: pitch -> architecture & concepts -> demo -> code walkthrough -> close.
Run: python presentation/build_deck_8min.py
  -> presentation/Omnichannel_8min.pptx  (+ presentation/slides/*.png previews)
"""
import os, math
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches

HERE = os.path.dirname(__file__)
SLIDES = os.path.join(HERE, "slides")
os.makedirs(SLIDES, exist_ok=True)

# ---- palette ----
RED=(196,18,48); INK=(15,23,42); SLATE=(30,41,59); SLATE2=(23,34,55)
CARD_D=(36,48,71); GREY=(82,92,107); MIST=(138,147,163); LINE=(214,221,230)
WHITE=(255,255,255); PAPER=(246,248,251); SHADOW=(228,232,238)
PINK_HI=(241,154,166); STEEL=(203,213,225); STEEL2=(194,204,217)
BLUE=(37,99,235); INDIGO=(79,70,229); GREEN=(21,128,61); AMBER=(180,84,9); TEAL=(15,118,110)
BLUE_T=(230,238,255); INDIGO_T=(237,236,254); GREEN_T=(230,244,234)
AMBER_T=(252,240,225); RED_T=(252,231,235); TEAL_T=(226,242,240)

W, H = 2560, 1440
FREG="/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf"
FBLD="/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf"
FIT ="/usr/share/fonts/truetype/liberation/LiberationSans-Italic.ttf"
FMONO="/usr/share/fonts/truetype/liberation/LiberationMono-Bold.ttf"
def F(sz, bold=False, ital=False, mono=False):
    return ImageFont.truetype(FMONO if mono else FBLD if bold else FIT if ital else FREG, sz)


class Canvas:
    def __init__(self, bg=WHITE):
        self.img = Image.new("RGB", (W, H), bg)
        self.d = ImageDraw.Draw(self.img)

    def tc(self, cx, cy, s, f, fill):
        l,t,r,b = self.d.textbbox((0,0), s, font=f)
        self.d.text((cx-(r-l)/2, cy-(b-t)/2-t), s, font=f, fill=fill)

    def tl(self, x, y, s, f, fill):
        self.d.text((x, y), s, font=f, fill=fill)

    def rrect(self, box, rad, fill, outline=None, width=3, shadow=False):
        if shadow:
            x0,y0,x1,y1=box; self.d.rounded_rectangle((x0+6,y0+8,x1+6,y1+8), rad, fill=SHADOW)
        self.d.rounded_rectangle(box, rad, fill=fill, outline=outline, width=width)

    def chip(self, x, y, label, tint, tc, fpx=30, padx=26, h=66, outline=True):
        f=F(fpx, bold=True)
        l,t,r,b=self.d.textbbox((0,0), label, font=f); w=(r-l)+padx*2
        self.rrect((x,y,x+w,y+h), h//2, tint, outline=(tc if outline else None), width=3)
        self.tc(x+w/2, y+h/2, label, f, tc)
        return w

    def save(self, name):
        p=os.path.join(SLIDES, name); self.img.save(p); return p


def header(c, kicker, title):
    c.d.rounded_rectangle((96,96,128,210), 6, fill=RED)              # red tick
    c.tl(150, 92, kicker.upper(), F(26, bold=True), RED)
    c.tl(148, 132, title, F(54, bold=True), INK)
    c.d.line((150, 252, 2410, 252), fill=LINE, width=3)


def footer(c, page, dark=False):
    col = (107,116,134) if dark else MIST
    c.tl(100, 1372, "Omnichannel Retail Analytics Assistant  ·  CMU Capstone", F(20), col)
    c.tc(2440, 1388, str(page), F(22, bold=True), col)


def star_(c, cx, cy, R, color):
    pts=[]
    for i in range(10):
        ang=-math.pi/2 + i*math.pi/5
        r=R if i%2==0 else R*0.42
        pts.append((cx+r*math.cos(ang), cy+r*math.sin(ang)))
    c.d.polygon(pts, fill=color)


def arrow_r(c, x0, x1, y, color=MIST, w=8):
    c.d.line((x0, y, x1-18, y), fill=color, width=w)
    c.d.polygon([(x1,y),(x1-24,y-15),(x1-24,y+15)], fill=color)

def arrow_u(c, x, y0, y1, color=LINE, w=9):
    c.d.line((x, y0, x, y1+18), fill=color, width=w)
    c.d.polygon([(x,y1),(x-15,y1+22),(x+15,y1+22)], fill=color)


# =================================================== diagram (reusable)
def make_arch(width=2360, height=1010):
    im = Image.new("RGB", (width, height), WHITE)
    cc = Canvas.__new__(Canvas); cc.img=im; cc.d=ImageDraw.Draw(im)
    m=24
    # concept chips
    cps=[("ReAct orchestration",BLUE_T,BLUE),("RAG",INDIGO_T,INDIGO),
         ("Knowledge graph",TEAL_T,TEAL),("Multi-agent",BLUE_T,BLUE),
         ("Tree-of-Thoughts + beam",RED_T,RED),("Guardrails",AMBER_T,AMBER),
         ("LLM grounding",GREEN_T,GREEN)]
    x=m
    for lab,ti,tcl in cps: x+=cc.chip(x,8,lab,ti,tcl,fpx=23,padx=20,h=52)+16
    # guardrails band
    cc.rrect((m,96,width-m,188),18,AMBER_T,outline=AMBER,width=3)
    cc.tc(width/2,142,"GUARDRAILS   ·   read-only SQL   ·   version / freshness gate   ·   "
          "write-refusal   ·   append-only audit trail",F(27,bold=True),AMBER)
    # flow
    flow=[("User question","natural language",BLUE_T,BLUE,None),
          ("Orchestrator","LangGraph ReAct loop",INDIGO_T,INDIGO,None),
          ("Multi-agent team","7 analysts · in parallel",BLUE_T,BLUE,None),
          ("Critic + beam","0–14 rubric → review",RED_T,RED,None),
          ("Governed answer","evidence · confidence · action",GREEN_T,GREEN,"LLM drafts wording only")]
    n=len(flow); gap=66; bw=((width-2*m)-gap*(n-1))/n; y0,y1=232,556; xs=[]
    for i,(t1,t2,ti,tcl,note) in enumerate(flow):
        x0=m+i*(bw+gap); x1=x0+bw; xs.append((x0,x1))
        cc.rrect((x0,y0,x1,y1),22,ti,outline=tcl,width=4,shadow=True)
        cx=(x0+x1)/2
        cc.tc(cx,y0+(y1-y0)*(0.34 if note else 0.40),t1,F(34,bold=True),tcl)
        cc.tc(cx,y0+(y1-y0)*(0.60 if note else 0.66),t2,F(23),GREY)
        if note: cc.tc(cx,y1-30,note,F(19,ital=True),GREEN)
    for i in range(n-1): arrow_r(cc,xs[i][1]+8,xs[i+1][0]-8,(y0+y1)/2)
    # foundation
    cc.tl(m,588,"GOVERNED FOUNDATION  (feeds the loop)",F(23,bold=True),GREY)
    fnd=[("YAML catalog","source of truth",RED_T,RED),("ChromaDB","RAG retrieval",INDIGO_T,INDIGO),
         ("NetworkX","knowledge graph",TEAL_T,TEAL),("DuckDB","medallion data · read-only",GREEN_T,GREEN)]
    fn=len(fnd); fgap=42; fbw=((width-2*m)-fgap*(fn-1))/fn; fy0,fy1=686,976
    for i,(t1,t2,ti,tcl) in enumerate(fnd):
        x0=m+i*(fbw+fgap); x1=x0+fbw
        arrow_u(cc,(x0+x1)/2,680,632)
        cc.rrect((x0,fy0,x1,fy1),22,ti,outline=tcl,width=4,shadow=True)
        cx=(x0+x1)/2
        cc.tc(cx,fy0+(fy1-fy0)*0.40,t1,F(34,bold=True),tcl)
        cc.tc(cx,fy0+(fy1-fy0)*0.66,t2,F(23),GREY)
    im.save(os.path.join(HERE,"architecture.png"))
    return im


# =================================================== SLIDE 1 — title
def slide1():
    c=Canvas(SLATE)
    c.d.rectangle((0,0,40,H), fill=RED)                         # left rail
    c.d.rectangle((1670,0,W,H), fill=SLATE2)                    # accent block
    c.tl(120,150,"CARNEGIE MELLON UNIVERSITY · FINAL CAPSTONE",F(26,bold=True),PINK_HI)
    c.tl(116,300,"Omnichannel Retail",F(96,bold=True),WHITE)
    c.tl(116,420,"Analytics Assistant",F(96,bold=True),WHITE)
    c.d.rectangle((122,580,460,592), fill=RED)
    c.tl(120,660,"A governed AI analyst that tells retail leaders why a KPI",F(40),STEEL)
    c.tl(120,720,"moved — across every team, in seconds, using only",F(40),STEEL)
    c.tl(120,780,"certified data. ",F(40),STEEL)
    w=c.d.textbbox((0,0),"certified data. ",font=F(40))[2]
    c.tl(120+w,780,"It never guesses.",F(40,bold=True),WHITE)
    ag=["1 · Pitch","2 · Architecture & concepts","3 · Live demo","4 · Code walkthrough"]
    x=120
    for lab in ag: x+=c.chip(x,1010,lab,CARD_D,WHITE,fpx=26,padx=26,h=70,outline=False)+30
    c.tl(120,1180,"Sreddy  ·  June 2026  ·  free & local, read-only prototype",F(26),MIST)
    return c.save("s1.png")


# =================================================== SLIDE 2 — architecture
def slide2():
    c=Canvas(WHITE)
    header(c,"Part 2 · Architecture & the concepts I used",
           "Governed knowledge + a multi-agent reasoning loop")
    arch=make_arch()
    scale=2360/2360  # native
    target_w=2360
    ar=arch.resize((target_w,int(arch.height*target_w/arch.width)))
    c.img.paste(ar,(int((W-target_w)/2),300))
    footer(c,2)
    return c.save("s2.png")


# =================================================== SLIDE 3 — demo
def slide3():
    c=Canvas(SLATE)
    c.d.rectangle((0,0,40,H), fill=RED)
    c.tl(120,150,"PART 3 · LIVE DEMO  ·  ~2.5 MIN",F(28,bold=True),PINK_HI)
    c.tl(116,210,"Seeing it run on a seeded “bad day”",F(64,bold=True),WHITE)
    rows=[("1","Why did conversion drop yesterday?",
           "Watch the team dispatch + live trace → drivers & evidence chart → owner actions + the human-review banner.",BLUE),
          ("2","“Update the paid-social budget”",
           "The read-only guardrail refuses the write — governance enforced in code, not a prompt.",AMBER),
          ("3","An executive briefing",
           "The same multi-agent path ranks the biggest cross-functional issues and routes each to an owner.",GREEN)]
    y=430
    for n,title,body,accent in rows:
        c.rrect((120,y,2440,y+250),26,CARD_D)
        c.d.rounded_rectangle((120,y,150,y+250),12,fill=accent)
        c.rrect((200,y+78,300,y+178),18,accent); c.tc(250,y+128,n,F(54,bold=True),WHITE)
        c.tl(350,y+58,title,F(40,bold=True),WHITE)
        # wrap body
        words=body.split(); line=""; yy=y+128
        for wd in words:
            test=(line+" "+wd).strip()
            if c.d.textbbox((0,0),test,font=F(28))[2] > 2000:
                c.tl(350,yy,line,F(28),STEEL); line=wd; yy+=44
            else: line=test
        c.tl(350,yy,line,F(28),STEEL)
        y+=290
    footer(c,3,dark=True)
    return c.save("s3.png")


# =================================================== SLIDE 4 — code walkthrough
def slide4():
    c=Canvas(WHITE)
    header(c,"Part 4 · Quick code walkthrough","Six files — each is one concept, end to end")
    rows=[("catalog/metrics.yaml","the certified definitions everything resolves to","governance / source of truth",RED_T,RED,True),
          ("sql_skill.py · check_sql()","only read-only SELECTs over approved tables run","guardrails",AMBER_T,AMBER,True),
          ("retrieval_skill.py · retrieve()","top-k vector lookup + a staleness / sync gate","RAG + embeddings",INDIGO_T,INDIGO,False),
          ("graph.py · build_workflow()","the LangGraph loop, routing & human-review gate","orchestration",BLUE_T,BLUE,True),
          ("agents/team.py · dispatch()","7 analysts in parallel; a failure is isolated","multi-agent systems",BLUE_T,BLUE,False),
          ("tot_skill.py · score_branch()","0–14 rubric; ungoverned hypothesis auto-pruned","Tree-of-Thoughts / beam",RED_T,RED,True)]
    y=300; rh=150
    for i,(mono,desc,concept,ti,tcl,star) in enumerate(rows):
        c.rrect((96,y,2410,y+rh-18),20,PAPER,outline=LINE,width=2)
        c.rrect((128,y+26,210,y+108),16,tcl); c.tc(169,y+67,str(i+1),F(40,bold=True),WHITE)
        c.tl(250,y+28,mono,F(33,mono=True),INK)
        if star:
            tw=c.d.textbbox((0,0),mono,font=F(33,mono=True))[2]
            star_(c,250+tw+36,y+46,17,(217,119,6))
        c.tl(250,y+86,desc,F(26),GREY)
        # concept chip right-aligned
        f=F(24,bold=True); tw=c.d.textbbox((0,0),concept,font=f)[2]; cw=tw+52
        c.rrect((2370-cw,y+40,2370,y+106),33,ti,outline=tcl,width=3); c.tc(2370-cw/2,y+73,concept,f,tcl)
        y+=rh
    star_(c,112,y+28,15,(217,119,6))
    c.tl(140,y+6,"keep these four if short on time.  Close on tot_skill.py: the ungoverned “maybe prices "
         "rose” idea is pruned automatically — governance isn't a prompt, it's enforced.",F(23,ital=True),GREY)
    footer(c,4)
    return c.save("s4.png")


# =================================================== SLIDE 5 — close
def slide5():
    c=Canvas(SLATE)
    c.d.rectangle((0,0,40,H), fill=RED)
    c.tl(120,210,"CLOSE",F(28,bold=True),PINK_HI)
    c.tl(116,300,"From a 2-day fire drill to a",F(72,bold=True),WHITE)
    c.tl(116,400,"governed answer in seconds.",F(72,bold=True),WHITE)
    c.d.rectangle((122,560,460,572), fill=RED)
    c.tl(120,640,"The hard part of agentic AI isn't generating answers —",F(38),STEEL)
    c.tl(120,700,"it's governing, evaluating, and observing them.",F(38,bold=True),WHITE)
    proof=["Seeded ground truth","15 validation checks","25 tests","0 UI exceptions","100% free / local"]
    x=120
    for lab in proof: x+=c.chip(x,860,lab,CARD_D,STEEL,fpx=26,padx=26,h=70,outline=False)+28
    c.tl(120,1030,"Extends to production: swap synthetic data for real tables behind the same catalog.",F(30),STEEL2)
    c.tl(120,1110,"Thank you — questions?",F(40,bold=True),WHITE)
    return c.save("s5.png")


# =================================================== assemble PPTX
NOTES = {
 1:"0:00–0:35. Say the pitch verbatim (big text), then read the 4 chips as your agenda. Energy up.",
 2:"0:35–2:35 (~2 min) — centerpiece. Trace it: a question comes in; the LangGraph ORCHESTRATOR runs a "
   "governed loop — RETRIEVES certified context (RAG/ChromaDB), VALIDATES vs the YAML catalog, RELATES "
   "metric→drivers via a NetworkX knowledge graph. It dispatches a TEAM of 7 analysts in PARALLEL; a "
   "CRITIC scores each on a 0–14 rubric, BEAM keeps the strongest, every recommendation goes to an owner "
   "for review. Two anchors: the whole loop runs inside GUARDRAILS (top band: read-only, version-gated, "
   "write-refusing, audited), fed by the GOVERNED FOUNDATION (YAML = source of truth, DuckDB read-only). "
   "Close: 'the LLM only drafts the wording — every number comes from governed SQL.' Point at each chip.",
 3:"2:35–5:00. App running and WARMED UP. (1) Run the conversion investigation; narrate the steps as the "
   "trace streams; land on Business answer → Evidence chart → owner actions + 🔴 human-review banner. "
   "(2) Type a write request → it refuses (fast, always lands). (3) One briefing. Short on time → #1 and "
   "#2 only. Keep a backup screenshot.",
 4:"5:00–7:30 (~2.5 min). Switch to VS Code with these SIX files PRE-OPENED as tabs. Open in order; ONE "
   "sentence each. Goal: show you understand the concepts, not read code. End on ungoverned_branch() being "
   "auto-pruned. Short on time → the four starred files. Bump editor font so it's readable on video.",
 5:"7:30–8:00. Tie back to fast AND trustworthy. Deliver the one-sentence lesson slowly. Sweep the proof "
   "chips in one breath. One line on the production path, thank them, open Q&A.",
}

def build():
    pngs=[slide1(), slide2(), slide3(), slide4(), slide5()]
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    blank=prs.slide_layouts[6]
    for i,png in enumerate(pngs,1):
        s=prs.slides.add_slide(blank)
        s.shapes.add_picture(png, 0, 0, width=prs.slide_width, height=prs.slide_height)
        s.notes_slide.notes_text_frame.text = NOTES[i]
    out=os.path.join(HERE,"Omnichannel_8min.pptx"); prs.save(out)
    print("Saved:", out, "| slides:", len(prs.slides._sldIdLst))

if __name__ == "__main__":
    build()
